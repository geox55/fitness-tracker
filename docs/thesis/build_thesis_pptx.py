"""Генератор презентаций к защите для Маши и Егора.

Запуск: `uv run --group ml python docs/thesis/build_thesis_pptx.py`

Артефакты: `docs/thesis/presentation-maria.pptx`, `presentation-egor.pptx`.

Структура слайдов (12 на каждого) описана в README; графики метрик
генерируются на лету через matplotlib и встраиваются как изображения.
"""

from __future__ import annotations

import io
from pathlib import Path

import matplotlib

matplotlib.use("Agg")  # без GUI-бекенда — для серверной/CI среды
import matplotlib.pyplot as plt

# ГОСТ: единый шрифт с засечками (Times New Roman) во всех графиках и схемах —
# как в тексте слайдов; рубленый DejaVu Sans по умолчанию выбивался из стиля.
plt.rcParams["font.family"] = "serif"
plt.rcParams["font.serif"] = ["Times New Roman", "DejaVu Serif"]
plt.rcParams["mathtext.fontset"] = "stix"
plt.rcParams["axes.unicode_minus"] = False

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

HERE = Path(__file__).parent
SCREENSHOTS = HERE.parent / "design" / "screenshots"
THESIS_SCREENSHOTS = HERE / "screenshots"
# Свежие скрины задеплоенного приложения (для слайда «Интерфейс приложения»).
STORE_SCREENSHOTS = HERE / "store" / "screenshots"

# --- Палитра -----------------------------------------------------------------
NAVY = RGBColor(0x14, 0x1B, 0x2D)        # тёмно-синий фон титула
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x0E, 0xA5, 0xE9)      # portal-blue — primary из дизайн-системы
TEXT_DARK = RGBColor(0x1F, 0x1F, 0x1F)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x55)
DIVIDER = RGBColor(0xE0, 0xE0, 0xE0)

# --- ЧБ-палитра для схем и графиков (только чёрный/белый, без серого) --------
# Различение — заливкой (чёрная заливка = акцент/ядро, белая с чёрной рамкой =
# обычный блок) и штриховкой (hatch) для серий на графиках, как в ГОСТ-графике.
BW_BLACK = "#000000"
BW_WHITE = "#FFFFFF"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# =============================================================================
# Низкоуровневые хелперы
# =============================================================================


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    """Чистый слайд без layout-плейсхолдеров — мы раскладываем сами."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_textbox(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int = 20,
    bold: bool = False,
    color: RGBColor = TEXT_DARK,
    align=PP_ALIGN.LEFT,
    font_name: str = "Times New Roman",
):
    """Универсальный текст-блок. Координаты в дюймах."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_bullets(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    items: list[str],
    font_size: int = 18,
    color: RGBColor = TEXT_DARK,
    space_after: int = 18,
    line_spacing: float = 1.15,
    vertical_center: bool = True,
):
    """Маркированный список. Каждый пункт — один параграф с буллетом «•».

    Блок вертикально центрируется в своей области (`vertical_center`) и
    раздаётся интервалами (`space_after`), чтобы 4-5 коротких пунктов
    заполняли слайд, а не липли к верхнему краю.
    """
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    if vertical_center:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = "Times New Roman"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def _add_rect(slide, *, left: float, top: float, width: float, height: float, color: RGBColor):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_accent_bar(slide, *, top: float = 0.65):
    """Тонкая фиолетовая полоса под заголовком — общий визуальный якорь."""
    _add_rect(slide, left=0.6, top=top, width=1.2, height=0.06, color=ACCENT)


def _header(slide, title: str):
    """Стандартный заголовок content-слайда (без подчёркивания — ГОСТ)."""
    _add_textbox(
        slide,
        left=0.6, top=0.3, width=12.0, height=0.6,
        text=title, font_size=28, bold=True, color=TEXT_DARK,
    )


def _footer(slide, page: int, total: int, author_short: str):
    """Номер слайда внизу справа."""
    _add_textbox(
        slide,
        left=11.3, top=6.85, width=1.5, height=0.5,
        text=f"{page}",
        font_size=20, color=TEXT_MUTED, align=PP_ALIGN.RIGHT,
    )


# =============================================================================
# matplotlib-чарты → PNG в memory
# =============================================================================


def _chart_inbody_models() -> io.BytesIO:
    """Сравнение MAE по 3 таргетам для 4 моделей (Маша)."""
    targets = ["Δ weight\n(кг)", "Δ body fat\n(%)", "Δ muscle\n(кг)"]
    persistence = [0.325, 0.276, 0.146]
    ridge = [0.168, 0.243, 0.127]
    lgbm = [0.156, 0.241, 0.120]
    mlp = [0.156, 0.238, 0.121]

    import numpy as np
    x = np.arange(len(targets))
    w = 0.2

    fig, ax = plt.subplots(figsize=(10, 5), dpi=160)
    # ЧБ: серии различаем заливкой и штриховкой, а не цветом.
    bars = [
        ("Persistence (baseline)", persistence, "white", "////"),
        ("Ridge", ridge, "white", "xxxx"),
        ("LightGBM (основная)", lgbm, "black", None),
        ("MLP (PyTorch)", mlp, "white", ".."),
    ]
    for i, (name, vals, fc, htch) in enumerate(bars):
        ax.bar(x + (i - 1.5) * w, vals, w, label=name, color=fc,
               edgecolor="black", linewidth=1.0, hatch=htch)

    ax.set_xticks(x)
    ax.set_xticklabels(targets, fontsize=11)
    ax.set_ylabel("MAE (меньше — лучше)", fontsize=11)
    # Легенду выносим НАД график (горизонтально, 4 в ряд), чтобы обозначения
    # не перекрывали столбцы.
    ax.legend(loc="lower center", bbox_to_anchor=(0.5, 1.01), ncol=4,
              fontsize=10, frameon=False)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", linestyle=":", color="black", linewidth=0.4)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_feature_importance_maria() -> io.BytesIO:
    """Feature importance для LGBM q50 weight (Маша) — горизонтальные бары.

    Цифры взяты из реального booster'а
    `ml/models/inbody_pred/lgbm/v0.1.0/delta_weight_kg_q50.joblib` (gain).
    """
    features = [
        "goal_weight_loss", "bmi", "goal_muscle_gain",
        "body_fat_percent", "weight_kg", "height_cm",
        "training_volume_t", "calories_t", "muscle_mass_kg",
        "ffm", "age", "sex_male",
    ]
    importance = [11488, 3162, 821, 278, 250, 247, 184, 167, 138, 132, 124, 15]
    # ЧБ: важные признаки — сплошной чёрный, остальные — белые со штриховкой.
    is_top = [v >= 1000 for v in importance]

    fig, ax = plt.subplots(figsize=(10, 5.2), dpi=160)
    y_pos = range(len(features))[::-1]  # сверху — самые важные
    bars = ax.barh(y_pos, importance,
                   color=["black" if t else "white" for t in is_top],
                   edgecolor="black", linewidth=1.0)
    for b, t in zip(bars, is_top):
        if not t:
            b.set_hatch("////")
    ax.set_yticks(y_pos)
    ax.set_yticklabels(features, fontsize=10)
    ax.set_xlabel("Importance (gain)", fontsize=11)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="x", linestyle=":", color="black", linewidth=0.4)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_feature_importance_egor() -> io.BytesIO:
    """Feature importance для LGBM workout-recommender (Егор), топ-10.

    Цифры из реального booster'а `ml/models/workout_rec/lgbm/v0.1.0/lgbm.joblib`.
    Модель доминируется user_equipment_count — это методологически
    осмысленно: рекомендовать упражнение со штангой пользователю без
    штанги бессмысленно.
    """
    features = [
        "user_equipment_count", "needs_bodyweight_only", "goal_weight_loss",
        "region_upper", "group_chest", "region_lower",
        "needs_dumbbell", "group_abs", "group_back", "needs_barbell",
    ]
    importance = [460204, 81609, 63268, 47979, 36857, 28712, 26758, 23041, 19961, 16042]
    # ЧБ: топ-признаки — сплошной чёрный, остальные — белые со штриховкой.
    is_top = [v >= 50000 for v in importance]

    fig, ax = plt.subplots(figsize=(10, 5.2), dpi=160)
    y_pos = range(len(features))[::-1]
    bars = ax.barh(y_pos, importance,
                   color=["black" if t else "white" for t in is_top],
                   edgecolor="black", linewidth=1.0)
    for b, t in zip(bars, is_top):
        if not t:
            b.set_hatch("////")
    ax.set_yticks(y_pos)
    ax.set_yticklabels(features, fontsize=10)
    ax.set_xlabel("Importance (gain)", fontsize=11)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="x", linestyle=":", color="black", linewidth=0.4)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_workout_models() -> io.BytesIO:
    """Сравнение ROC-AUC для popularity / LR / LGBM (Егор)."""
    models = ["Popularity\n(baseline)", "Logistic\nRegression", "LightGBM\n(основная)"]
    roc_auc = [0.500, 0.939, 0.985]

    fig, ax = plt.subplots(figsize=(10, 5), dpi=160)
    # ЧБ: baseline и LR — белые со штриховкой, основная модель — сплошной чёрный.
    bars = ax.bar(models, roc_auc, color=["white", "white", "black"],
                  edgecolor="black", width=0.55, linewidth=1.0)
    bars[0].set_hatch("////")
    bars[1].set_hatch("xxxx")
    ax.set_ylim(0.4, 1.05)
    ax.set_ylabel("ROC-AUC (больше — лучше)", fontsize=11)
    for b, v in zip(bars, roc_auc):
        ax.text(b.get_x() + b.get_width() / 2, v + 0.015, f"{v:.3f}",
                ha="center", fontsize=12, fontweight="bold")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", linestyle=":", color="black", linewidth=0.4)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _diagram_architecture() -> io.BytesIO:
    """Схема архитектуры приложения Portal — блоки со стрелками."""
    fig, ax = plt.subplots(figsize=(11, 5.5), dpi=160)
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 6)
    ax.axis("off")

    from matplotlib.patches import FancyBboxPatch, FancyArrowPatch

    # ЧБ: акцентные блоки (узел API и ML-модель) — чёрные с белым текстом,
    # остальные — белые с чёрной рамкой и чёрным текстом.
    def box(x, y, w, h, label, fill="white", text_color="black", fontsize=12):
        rect = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.15",
                              facecolor=fill, edgecolor="black", linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x + w / 2, y + h / 2, label, ha="center", va="center",
                fontsize=fontsize, fontweight="bold", color=text_color)

    def arrow(x1, y1, x2, y2, label=""):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="->", color="black", lw=2))
        if label:
            mx, my = (x1 + x2) / 2, (y1 + y2) / 2
            ax.text(mx + 0.15, my, label, fontsize=9, color="black")

    box(3.5, 4.8, 4.0, 0.9, "Flutter PWA\n(браузер / мобильный)")
    box(3.5, 2.8, 4.0, 0.9, "FastAPI\n(REST API)", "black", "white")
    box(0.3, 0.5, 2.5, 0.9, "PostgreSQL 16", fontsize=11)
    box(4.2, 0.5, 2.5, 0.9, "MinIO\n(файлы, PDF)")
    box(8.2, 0.5, 2.5, 0.9, "LightGBM\n(ML-прогноз)", "black", "white")
    box(8.2, 2.8, 2.5, 0.9, "Watcher\n(адаптация)")

    # Концы стрелок отводим к ВИДИМЫМ краям блоков с зазором: boxstyle pad=0.15
    # раздувает рамку на 0.15 за пределы rect, поэтому стрелки к номинальным
    # краям (4.8/3.7/2.8/1.4) втыкались внутрь карточек.
    arrow(5.5, 4.6, 5.5, 3.9, "/api/v1/*")   # PWA → FastAPI
    arrow(4.6, 2.6, 1.7, 1.62)               # FastAPI → PostgreSQL
    arrow(5.5, 2.6, 5.45, 1.62)              # FastAPI → MinIO
    arrow(6.4, 2.6, 9.3, 1.62)               # FastAPI → LightGBM
    arrow(7.7, 3.25, 8.0, 3.25)              # FastAPI → Watcher

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def _diagram_forecast_pipeline(show_note: bool = True) -> io.BytesIO:
    """Пайплайн прогноза InBody — ровная горизонтальная цепочка.

    Шесть одинаковых блоков, выровненных по центру с равными промежутками;
    пояснение про вход и cold start — отдельной строкой по центру снизу
    (а не косыми выносками сбоку, из-за которых картинка съезжала и
    обрезалась неровно при tight-кадрировании)."""
    fig, ax = plt.subplots(figsize=(12, 2.7), dpi=160)
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 2.7)
    ax.axis("off")

    from matplotlib.patches import FancyBboxPatch

    # ЧБ: ключевой блок (ML-модель) — чёрный с белым текстом, остальные —
    # белые с чёрной рамкой и чёрным текстом.
    steps = [
        ("InBody\nзамер", "white", "black"),
        ("12\nпризнаков", "white", "black"),
        ("LightGBM\nq10 / q50 / q90", "black", "white"),
        ("Прогноз Δ\n(вес, жир, мышцы)", "white", "black"),
        ("Интервал 80 %\n+ пояснение", "white", "black"),
        ("Советы\n(ВОЗ / ACSM)", "white", "black"),
    ]

    box_w, box_h, gap = 1.7, 1.2, 0.3
    n = len(steps)
    total = n * box_w + (n - 1) * gap
    start = (12 - total) / 2          # симметричные поля слева/справа
    y = 0.95
    cy = y + box_h / 2

    centers = []
    for i, (label, fill, text_color) in enumerate(steps):
        x = start + i * (box_w + gap)
        rect = FancyBboxPatch((x, y), box_w, box_h,
                              boxstyle="round,pad=0.03,rounding_size=0.10",
                              facecolor=fill, edgecolor="black", linewidth=1.4)
        ax.add_patch(rect)
        ax.text(x + box_w / 2, cy, label, ha="center", va="center",
                fontsize=12, fontweight="bold", color=text_color)
        centers.append((x, x + box_w))

    for i in range(n - 1):
        ax.annotate("", xy=(centers[i + 1][0], cy), xytext=(centers[i][1], cy),
                    arrowprops=dict(arrowstyle="-|>", color="black", lw=2.2))

    if show_note:
        ax.text(6.0, 0.4,
                "Вход — профиль пользователя и история замеров; "
                "при < 4 замерах применяется baseline-прогноз (REQ-12)",
                ha="center", va="center", fontsize=10, style="italic",
                color="black")

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_ci_coverage() -> io.BytesIO:
    """Калибровка доверительных интервалов: фактический coverage vs теоретический 80%."""
    targets = ["Δ weight", "Δ body fat", "Δ muscle"]
    lgbm = [0.79, 0.78, 0.80]
    mlp = [0.83, 0.79, 0.83]
    target_line = 0.80

    import numpy as np
    x = np.arange(len(targets))
    w = 0.35

    fig, ax = plt.subplots(figsize=(9, 4.5), dpi=160)
    # ЧБ: LightGBM — сплошной чёрный, MLP — белый со штриховкой.
    ax.bar(x - w / 2, lgbm, w, label="LightGBM", color="black", edgecolor="black",
           linewidth=1.0)
    b_mlp = ax.bar(x + w / 2, mlp, w, label="MLP", color="white", edgecolor="black",
                   linewidth=1.0)
    for b in b_mlp:
        b.set_hatch("////")
    ax.axhline(y=target_line, color="black", linestyle="--", linewidth=1.5,
               label="Теоретический уровень 80%")
    ax.set_xticks(x)
    ax.set_xticklabels(targets, fontsize=11)
    ax.set_ylabel("Доля точек в [q10, q90]", fontsize=11)
    ax.set_ylim(0.60, 0.95)
    ax.legend(loc="lower right", fontsize=10, frameon=False)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", linestyle=":", color="black", linewidth=0.4)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


# =============================================================================
# Слайд-конструкторы
# =============================================================================


def slide_title(prs: Presentation, *, topic: str, author: str, group: str,
                supervisor: str = "_______________",
                program: str | None = None, student_word: str = "Студент"):
    """Титульный слайд по образцу титульного листа диплома (тот же текст и
    порядок блоков), но БЕЗ подписей и грифа «УТВЕРЖДАЮ» — на слайде они не
    нужны. Всё по центру, тема — крупно посередине, руководитель и студент —
    внизу справа, как на листе."""
    s = _blank_slide(prs)
    C = PP_ALIGN.CENTER
    R = PP_ALIGN.RIGHT

    # Шапка вуза — по центру.
    _add_textbox(s, left=0.6, top=0.35, width=12.13, height=0.35,
                 text="МИНОБРНАУКИ РОССИИ", font_size=14, bold=True,
                 color=TEXT_DARK, align=C)
    _add_textbox(s, left=0.9, top=0.7, width=11.53, height=0.85,
                 text="федеральное государственное автономное образовательное "
                      "учреждение высшего образования «Омский государственный "
                      "университет им. Ф.М. Достоевского»",
                 font_size=12, color=TEXT_DARK, align=C)
    _add_textbox(s, left=0.6, top=1.5, width=12.13, height=0.35,
                 text="Кафедра компьютерной математики и программного обеспечения",
                 font_size=12, color=TEXT_MUTED, align=C)

    # Тема — крупно по центру.
    _add_textbox(s, left=0.8, top=2.6, width=11.73, height=1.5,
                 text=topic, font_size=26, bold=True, color=TEXT_DARK, align=C)

    # Вид работы и направление.
    _add_textbox(s, left=0.6, top=4.2, width=12.13, height=0.35,
                 text="Выпускная квалификационная работа", font_size=14,
                 bold=True, color=TEXT_DARK, align=C)
    _add_textbox(s, left=0.6, top=4.55, width=12.13, height=0.35,
                 text="по направлению 01.04.02 — «Прикладная математика "
                      "и информатика»",
                 font_size=12, color=TEXT_MUTED, align=C)
    if program:
        _add_textbox(s, left=0.6, top=4.87, width=12.13, height=0.35,
                     text=f"магистерская программа «{program}»",
                     font_size=12, color=TEXT_MUTED, align=C)

    # Руководитель и студент — внизу справа (как на листе), без строк подписи.
    sup = supervisor.replace("\n", " ")
    _add_textbox(s, left=6.0, top=5.6, width=6.7, height=0.32,
                 text="Научный руководитель:", font_size=12, color=TEXT_MUTED,
                 align=R)
    _add_textbox(s, left=5.5, top=5.9, width=7.2, height=0.55,
                 text=sup, font_size=13, color=TEXT_DARK, align=R)
    _add_textbox(s, left=6.0, top=6.45, width=6.7, height=0.32,
                 text=f"{student_word} гр. {group}:", font_size=12, color=TEXT_MUTED,
                 align=R)
    _add_textbox(s, left=6.0, top=6.73, width=6.7, height=0.32,
                 text=author, font_size=13, bold=True, color=TEXT_DARK, align=R)

    # Город и год — по центру внизу.
    _add_textbox(s, left=0.6, top=7.05, width=12.13, height=0.35,
                 text="Омск, 2026", font_size=12, color=TEXT_MUTED, align=C)


def slide_relevance(prs, *, page, total, author_short, intro: str, points: list[str]):
    s = _blank_slide(prs)
    _header(s, "Актуальность")
    _add_textbox(s, left=0.6, top=1.1, width=12.0, height=0.8,
                 text=intro, font_size=24, color=TEXT_DARK)
    _add_bullets(s, left=0.6, top=2.3, width=12.0, height=4.7, items=points,
                 font_size=26)
    _footer(s, page, total, author_short)


def slide_goal_tasks(prs, *, page, total, author_short, goal: str, tasks: list[str]):
    s = _blank_slide(prs)
    _header(s, "Цель и задачи")
    _add_textbox(s, left=0.6, top=1.1, width=2.5, height=0.4,
                 text="Цель", font_size=22, bold=True, color=TEXT_DARK)
    _add_textbox(s, left=0.6, top=1.5, width=12.0, height=1.2,
                 text=goal, font_size=24, color=TEXT_DARK)
    _add_textbox(s, left=0.6, top=2.9, width=4.0, height=0.4,
                 text="Задачи", font_size=22, bold=True, color=TEXT_DARK)
    _add_bullets(s, left=0.6, top=3.3, width=12.0, height=3.7, items=tasks,
                 font_size=24)
    _footer(s, page, total, author_short)


def slide_architecture(prs, *, page, total, author_short, points: list[str], stack: list[str]):
    s = _blank_slide(prs)
    _header(s, "Архитектура приложения")
    _add_textbox(s, left=0.6, top=1.1, width=6.0, height=0.4,
                 text="Принципы", font_size=22, bold=True, color=TEXT_DARK)
    _add_bullets(s, left=0.6, top=1.5, width=6.0, height=5.0, items=points,
                 font_size=24)
    _add_textbox(s, left=7.0, top=1.1, width=6.0, height=0.4,
                 text="Стек", font_size=22, bold=True, color=TEXT_DARK)
    _add_bullets(s, left=7.0, top=1.5, width=6.0, height=5.0, items=stack,
                 font_size=24)
    _footer(s, page, total, author_short)


# Сквозная нумерация таблиц в презентации (сбрасывается в начале build_*).
_PPTX_TBL = {"n": 0}


def slide_table(prs, *, page, total, author_short, title: str,
                headers: list[str], rows: list[list[str]],
                highlight_row: int | None = None, caption: str | None = None,
                table_name: str | None = None,
                col_widths: list[float] | None = None,
                table_w: float = 11.5):
    """Слайд с таблицей: заголовок + название «Таблица N. …» + таблица + подпись.

    `highlight_row` — индекс строки данных (0..len(rows)-1) для подсветки
    основной модели чёрным. Если None — без подсветки.
    `table_name` — название для подписи «Таблица N. …» над таблицей (ГОСТ);
    если None, берётся заголовок слайда.
    `col_widths` — ширины столбцов в дюймах (сумма = table_w); если None,
    столбцы равной ширины. Нужно, когда первый столбец короткий (например,
    ФИО), а второй длинный — иначе таблица смотрится несимметрично.
    """
    s = _blank_slide(prs)
    _header(s, title)

    n_rows = len(rows) + 1  # +header
    n_cols = len(headers)
    # Центрируем таблицу по ширине слайда (table_w задаётся параметром)
    header_h = 0.65
    row_h = 0.85
    table_h = header_h + row_h * len(rows)
    left = (13.333 - table_w) / 2
    # Вертикально центрируем блок таблицы в свободной зоне (под заголовком и
    # подписью «Таблица N», с запасом под нижнюю подпись/футер).
    zone_top, zone_bottom = 2.0, 6.7
    top = max(zone_top, zone_top + (zone_bottom - zone_top - table_h) / 2)

    # ГОСТ-подпись «Таблица N. Название» — справа над таблицей.
    _PPTX_TBL["n"] += 1
    _add_textbox(s, left=left, top=top - 0.5, width=table_w, height=0.4,
                 text=f"Таблица {_PPTX_TBL['n']}. {table_name or title}",
                 font_size=14, color=TEXT_DARK, align=PP_ALIGN.RIGHT)

    tbl_shape = s.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                   Inches(table_w), Inches(table_h))
    tbl = tbl_shape.table
    # ЧБ: стиль «Table Grid» — чёрная сетка, без цветных полос/заливок, которые
    # PowerPoint навешивает по умолчанию. Спец-форматирование первой строки и
    # чередование строк выключаем — оформляем ячейки вручную.
    tblPr = tbl._tbl.tblPr
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    style_id = tblPr.find(qn("a:tableStyleId"))
    if style_id is None:
        style_id = tblPr.makeelement(qn("a:tableStyleId"), {})
        tblPr.append(style_id)
    style_id.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # Table Grid
    # Явные высоты строк — иначе python-pptx раскидает их слишком низкими.
    tbl.rows[0].height = Inches(header_h)
    for i in range(1, n_rows):
        tbl.rows[i].height = Inches(row_h)
    # Ширины столбцов: по умолчанию равные, иначе — заданные (для таблиц с
    # коротким первым столбцом, чтобы содержимое выглядело симметрично).
    if col_widths:
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = Inches(w)

    # Заголовок: чёрный фон, белый текст
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.name = "Times New Roman"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = WHITE

    # Данные
    for i, row in enumerate(rows):
        is_highlight = highlight_row is not None and i == highlight_row
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # ЧБ: выделенная строка — инверсия (чёрная заливка, белый текст),
            # остальные — белые с чёрным текстом. Без серого.
            cell.fill.solid()
            cell.fill.fore_color.rgb = (RGBColor(0x00, 0x00, 0x00) if is_highlight
                                        else RGBColor(0xFF, 0xFF, 0xFF))
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            # Первая колонка — слева, остальные — по центру (для чисел красиво)
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.name = "Times New Roman"
            run.font.size = Pt(16)
            run.font.bold = is_highlight
            run.font.color.rgb = WHITE if is_highlight else TEXT_DARK

    if caption:
        _add_textbox(s, left=0.6, top=top + table_h + 0.3, width=12.0, height=0.7,
                     text=caption, font_size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    _footer(s, page, total, author_short)


def slide_content(prs, *, page, total, author_short, title: str, items: list[str],
                  intro: str | None = None):
    s = _blank_slide(prs)
    _header(s, title)
    top = 1.1
    if intro:
        _add_textbox(s, left=0.6, top=top, width=12.0, height=0.8,
                     text=intro, font_size=24, color=TEXT_DARK)
        top = 2.2
    _add_bullets(s, left=0.6, top=top, width=12.0, height=7.5 - top - 0.5,
                 items=items, font_size=26)
    _footer(s, page, total, author_short)


def slide_chart(prs, *, page, total, author_short, title: str, chart: io.BytesIO,
                caption: str | None = None):
    s = _blank_slide(prs)
    _header(s, title)
    # Центрируем по ШИРИНЕ (фикс. ширина) и по ВЫСОТЕ в свободной зоне под
    # заголовком — иначе широкие невысокие картинки (пайплайн, графики) липнут
    # к верху и оставляют большой пустой низ → выглядят «не по центру».
    pic = s.shapes.add_picture(chart, Inches(0), Inches(1.1),
                               width=Inches(11.0))
    zone_top = 1.05
    zone_bottom = 6.25 if caption else 6.7
    zone_h = zone_bottom - zone_top
    # Если картинка выше зоны — ужимаем с сохранением пропорций, чтобы не лезла
    # на подпись/футер. Иначе оставляем ширину 11".
    if pic.height.inches > zone_h:
        ratio = zone_h / pic.height.inches
        pic.width = Inches(pic.width.inches * ratio)
        pic.height = Inches(zone_h)
    pic.left = Inches((13.333 - pic.width.inches) / 2)
    pic.top = Inches(zone_top + (zone_h - pic.height.inches) / 2)
    if caption:
        _add_textbox(s, left=0.6, top=6.4, width=12.0, height=0.5,
                     text=caption, font_size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    _footer(s, page, total, author_short)


def slide_diagram(prs, *, page, total, author_short, title: str,
                  img_path, caption: str | None = None, img_h: float = 5.0):
    """Слайд с готовой картинкой (ER/UML): вписываем по высоте и центрируем
    по ширине — диаграммы шире обычных графиков, фиксированная ширина не
    подходит. add_picture с одной высотой сохраняет пропорции."""
    s = _blank_slide(prs)
    _header(s, title)
    top = 1.25
    try:
        pic = s.shapes.add_picture(str(img_path), Inches(0), Inches(top),
                                   height=Inches(img_h))
        pic.left = Inches((13.333 - pic.width.inches) / 2)
    except Exception:
        _add_rect(s, left=1.2, top=top, width=11.0, height=img_h, color=DIVIDER)
    if caption:
        _add_textbox(s, left=0.6, top=top + img_h + 0.15, width=12.0,
                     height=0.5, text=caption, font_size=13,
                     color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    _footer(s, page, total, author_short)


def slide_screenshots(prs, *, page, total, author_short, title: str,
                      shots: list[tuple[Path, str]], keep_aspect: bool = False):
    """До 3 скриншотов в ряд с подписями.

    keep_aspect=True вписывает каждый скрин в слот (each_w × each_h) с
    сохранением пропорций и центрированием — нужно для портретных мобильных
    скриншотов, которые иначе растягиваются по ширине. По умолчанию False,
    чтобы не менять уже свёрстанные слайды (дека Егора)."""
    s = _blank_slide(prs)
    _header(s, title)
    n = len(shots)
    total_w = 11.5
    gap = 0.3
    each_w = (total_w - gap * (n - 1)) / n
    each_h = 4.8
    top = 1.2
    start_left = (13.333 - total_w) / 2

    for i, (path, caption) in enumerate(shots):
        left = start_left + i * (each_w + gap)
        try:
            if keep_aspect:
                with PILImage.open(path) as im:
                    iw, ih = im.size
                # вписываем в слот, центрируем по горизонтали, выравниваем по верху
                scale = min(each_w / iw, each_h / ih)
                w, h = iw * scale, ih * scale
                s.shapes.add_picture(
                    str(path), Inches(left + (each_w - w) / 2), Inches(top),
                    width=Inches(w), height=Inches(h))
            else:
                s.shapes.add_picture(str(path), Inches(left), Inches(top),
                                     width=Inches(each_w), height=Inches(each_h))
        except Exception:
            _add_rect(s, left=left, top=top, width=each_w, height=each_h, color=DIVIDER)
        _add_textbox(s, left=left, top=6.15, width=each_w, height=0.6,
                     text=caption, font_size=12, color=TEXT_MUTED,
                     align=PP_ALIGN.CENTER)
    _footer(s, page, total, author_short)


def slide_conclusion(prs, *, page, total, author_short, items: list[str]):
    s = _blank_slide(prs)
    _header(s, "Выводы")
    _add_bullets(s, left=0.6, top=1.2, width=12.0, height=5.6, items=items,
                 font_size=28)
    _footer(s, page, total, author_short)


def slide_thanks(prs, *, author: str, qr_path=None):
    """Финальный слайд — белый фон (ГОСТ), «Спасибо», ФИО.

    qr_path — опциональный QR-код на развёрнутое приложение (чтобы комиссия
    отсканировала и попробовала вживую во время вопросов). Ссылку текстом не
    выводим — только QR. Если не задан — слайд как раньше (вариант Егора).
    """
    s = _blank_slide(prs)
    has_qr = qr_path is not None and Path(qr_path).exists()

    if not has_qr:
        _add_textbox(s, left=0.6, top=2.5, width=12.0, height=1.5,
                     text="Спасибо за внимание",
                     font_size=54, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        _add_textbox(s, left=0.6, top=4.4, width=12.0, height=0.6,
                     text="Готова ответить на вопросы",
                     font_size=24, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        _add_textbox(s, left=0.6, top=5.5, width=12.0, height=0.5,
                     text=author, font_size=16, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        _add_textbox(s, left=0.6, top=5.95, width=12.0, height=0.4,
                     text="Омск · 2026", font_size=12,
                     color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        return

    # Вариант с QR-кодом (демонстрация приложения).
    _add_textbox(s, left=0.6, top=0.9, width=12.0, height=1.2,
                 text="Спасибо за внимание",
                 font_size=46, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    # Задаём только высоту — ширину python-pptx считает из родных пропорций
    # картинки (QR с подписью «Scan Me» портретный 3000×3889), иначе при
    # width==height его сплющивает. Центрируем по фактической ширине.
    qr_h = 3.6
    qr_top = 2.15
    try:
        pic = s.shapes.add_picture(str(qr_path), Inches(0), Inches(qr_top),
                                   height=Inches(qr_h))
        pic.left = Inches((13.333 - pic.width.inches) / 2)
    except Exception:
        _add_rect(s, left=(13.333 - 2.8) / 2, top=qr_top, width=2.8,
                  height=qr_h, color=DIVIDER)


# =============================================================================
# СБОРКА: Маша (12 слайдов)
# =============================================================================


def build_maria() -> Path:
    _PPTX_TBL["n"] = 0
    prs = _new_prs()
    total = 17
    short = "Лапова М. С."

    # 1. Титул
    slide_title(
        prs,
        topic="РАЗРАБОТКА КРОССПЛАТФОРМЕННОГО ПРИЛОЖЕНИЯ "
              "ДЛЯ ПРОГНОЗА ПОКАЗАТЕЛЕЙ INBODY С АДАПТАЦИЕЙ ТРЕНИРОВОК",
        author="Лапова Мария Сергеевна",
        group="МММ-401-О-03",
        supervisor="Агафонов А. Л.,\nдоцент, каф. компьютерной математики "
                   "и программного обеспечения",
        student_word="Студентка",
    )

    # 2. Актуальность
    slide_content(
        prs, page=2, total=total, author_short=short,
        title="Актуальность",
        items=[
            "Рост популярности фитнес-приложений",
            "Отсутствие персонализированного прогноза",
            "Нет адаптации плана тренировок",
            "InBody используется без аналитики",
        ],
    )

    # 3. Цель и задачи
    slide_goal_tasks(
        prs, page=3, total=total, author_short=short,
        goal="Разработка приложения Portal для прогноза "
             "состава тела и адаптации тренировочного плана.",
        tasks=[
            "Анализ предметной области",
            "Разработка ML-модели прогноза",
            "Создание REST API",
            "Реализация PWA-приложения",
            "Тестирование системы",
        ],
    )

    # 4. Разделение работ (план работ — кто что делал)
    slide_table(
        prs, page=4, total=total, author_short=short,
        title="Разделение работ",
        headers=["Студент", "Что разрабатывал"],
        rows=[
            ["Лапова М. С.", "Прогноз состава тела (InBody), доверительные "
                             "интервалы, адаптация плана, синтез данных, "
                             "импорт PDF InBody"],
            ["Ощепков Е. С.", "Генератор плана тренировок, каталог "
                              "упражнений, учёт тренировок"],
            ["Совместно", "Концепция, архитектура, REST API, "
                          "PWA-клиент, база данных"],
        ],
        table_w=10.4,
        col_widths=[3.0, 7.4],
    )

    # 5. Сравнение с аналогами
    slide_table(
        prs, page=5, total=total, author_short=short,
        title="Сравнение с существующими решениями",
        headers=["Приложение", "Прогноз", "InBody", "Адаптация"],
        rows=[
            ["Hevy, Strong", "—", "—", "—"],
            ["Fitbod", "—", "—", "—"],
            ["MacroFactor", "Только вес", "—", "Частично"],
            ["DDX Fitness", "—", "Замер", "—"],
            ["Portal", "+", "+", "+"],
        ],
        highlight_row=4,
    )

    # 6. Архитектура (схема)
    slide_chart(
        prs, page=6, total=total, author_short=short,
        title="Архитектура системы",
        chart=_diagram_architecture(),
        caption="Рис. 1. Архитектура приложения Portal",
    )

    # 7. Схема базы данных (ER-диаграмма)
    slide_diagram(
        prs, page=7, total=total, author_short=short,
        title="Схема базы данных",
        img_path=HERE / "er-diagram-maria.png",
        caption="Рис. 2. Сущности и связи подсистемы прогноза и адаптации",
    )

    # 8. Основные функции
    slide_content(
        prs, page=8, total=total, author_short=short,
        title="Основные функции",
        items=[
            "Прогноз состава тела (вес, жир, мышцы)",
            "Генерация плана тренировок на 4 недели",
            "Импорт PDF-отчётов InBody",
            "Рекомендации по нормам ВОЗ и ACSM",
            "Автоадаптация плана при изменениях профиля",
        ],
    )

    # 9. ML-модель
    slide_content(
        prs, page=9, total=total, author_short=short,
        title="ML-модель прогноза",
        intro="Прогнозируются: вес, процент жира, мышечная масса "
              "на горизонтах 1, 2, 4 недели.",
        items=[
            "Ridge Regression — линейная базовая модель",
            "LightGBM — градиентный бустинг (основная модель)",
            "Нейросеть (PyTorch) — альтернативная архитектура",
            "Квантильная регрессия для доверительных интервалов",
        ],
    )

    # 10. Как работает прогноз (пайплайн модели)
    slide_chart(
        prs, page=10, total=total, author_short=short,
        title="Как работает прогноз",
        chart=_diagram_forecast_pipeline(show_note=False),
        caption="Рис. 3. Пайплайн прогноза InBody",
    )

    # 11. Последовательность обработки запроса (UML)
    slide_diagram(
        prs, page=11, total=total, author_short=short,
        title="Последовательность обработки запроса",
        img_path=HERE / "uml-sequence-maria.png",
        caption="Рис. 4. Диаграмма последовательности: запрос прогноза InBody",
    )

    # 12. Доверительные интервалы
    slide_content(
        prs, page=12, total=total, author_short=short,
        title="Доверительные интервалы",
        items=[
            "Квантильная регрессия: квантили 0,1 / 0,5 / 0,9",
            "Интервал 80%: факт попадает в 79% случаев",
            "L(y, ŷ) = max(τ·(y−ŷ), (τ−1)·(y−ŷ))",
            "9 бустеров: 3 метрики × 3 квантиля",
        ],
    )

    # 13. Алгоритм адаптации
    slide_content(
        prs, page=13, total=total, author_short=short,
        title="Алгоритм адаптации плана",
        intro="Триггеры обновления плана:",
        items=[
            "Изменение веса более чем на 2 кг",
            "Изменение цели тренировок",
            "Изменение частоты тренировок",
            "Завершение 4-недельного цикла",
        ],
    )

    # 14. Сравнение ошибок по моделям (таблица MAE)
    slide_table(
        prs, page=14, total=total, author_short=short,
        title="Сравнение ошибок по моделям",
        headers=["Модель", "MAE, кг", "RMSE, кг", "R²", "Лучше baseline на"],
        rows=[
            ["Persistence (baseline)", "0,325", "0,399", "−0,13", "—"],
            ["Ridge", "0,168", "0,210", "0,69", "48 %"],
            ["LightGBM (основная)", "0,156", "0,196", "0,73", "52 %"],
            ["Нейросеть (PyTorch)", "0,156", "0,196", "0,73", "52 %"],
        ],
        highlight_row=2,
    )

    # 15. Результаты (график)
    slide_chart(
        prs, page=15, total=total, author_short=short,
        title="Результаты: сравнение моделей",
        chart=_chart_inbody_models(),
        caption="Рис. 5. MAE четырёх моделей на тестовой выборке",
    )

    # 16. Интерфейс приложения (3 экрана задеплоенного Portal)
    slide_screenshots(
        prs, page=16, total=total, author_short=short,
        title="Интерфейс приложения",
        shots=[
            (STORE_SCREENSHOTS / "03_home.jpg",
             "Главный экран"),
            (STORE_SCREENSHOTS / "09_inbody_pdf.jpg",
             "Загрузка и анализ InBody"),
            (STORE_SCREENSHOTS / "08_statistics.jpg",
             "Статистика прогресса"),
        ],
        keep_aspect=True,
    )

    # 17. Заключение
    slide_conclusion(
        prs, page=17, total=total, author_short=short,
        items=[
            "Создано приложение Portal (FastAPI + Flutter)",
            "Реализован прогноз состава тела (R² = 0,73)",
            "Реализована автоадаптация тренировок",
            "Подтверждена работоспособность системы",
        ],
    )

    # 16. Спасибо (+ QR на развёрнутое приложение для демо)
    slide_thanks(prs, author="Лапова Мария Сергеевна · группа МММ-401-О-03",
                 qr_path=HERE / "qr-app.png")

    out = HERE / "presentation-maria.pptx"
    prs.save(out)
    return out


# =============================================================================
# СБОРКА: Егор (12 слайдов)
# =============================================================================


def build_egor() -> Path:
    _PPTX_TBL["n"] = 0
    prs = _new_prs()
    total = 14
    short = "Ощепков Е. С."

    # 1. Титул
    slide_title(
        prs,
        topic="Разработка кроссплатформенного приложения "
              "для автоматической генерации планов тренировок "
              "с использованием машинного обучения",
        author="Ощепков Егор Сергеевич",
        group="МММ-401-О-03",
        supervisor="Агафонов А. Л.,\nдоцент, каф. компьютерной математики "
                   "и программного обеспечения",
    )

    # 2. Актуальность
    slide_relevance(
        prs, page=2, total=total, author_short=short,
        intro="Самостоятельные тренировки без программы — основная причина "
              "плато и отсутствия прогресса у новичков и любителей.",
        points=[
            "Существующие фитнес-приложения либо ведут дневник без генерации "
            "программ (Hevy, Strong), либо предлагают шаблонные планы (Fitbod, Caliber)",
            "Универсальные программы не учитывают цель, уровень, частоту, доступное "
            "оборудование и динамику пользователя",
            "Чистый rule-based — однообразен; чистое ML — требует десятков тысяч пар "
            "(профиль, программа), которых нет в открытом доступе",
            "Решение: гибридный алгоритм — ML-ранкер упражнений + детерминированный "
            "composer, гарантирующий корректность сборки плана",
        ],
    )

    # 3. Цель и задачи
    slide_goal_tasks(
        prs, page=3, total=total, author_short=short,
        goal="Разработать кроссплатформенное приложение для автоматической "
             "генерации индивидуальных тренировочных программ на основе ML.",
        tasks=[
            "Проанализировать существующие подходы к генерации тренировочных "
            "программ (rule-based, end-to-end ML, рекомендательные системы)",
            "Сформировать обучающий датасет из 97 760 пар (профиль, упражнение, "
            "метка) на основе Kaggle-источника и rule-based разметки",
            "Обучить и сравнить три модели ранжирования: popularity-baseline, "
            "Logistic Regression, LightGBM Classifier",
            "Реализовать composer — детерминированный сборщик программы из "
            "ранжированного каталога с проверкой тренировочных принципов",
            "Спроектировать REST-сервис генерации с lazy-load ML-артефакта и "
            "rule-based fallback при отсутствии модели",
            "Покрыть критические компоненты автотестами; обеспечить версионирование "
            "моделей и воспроизводимость пайплайна",
        ],
    )

    # 4. Архитектура
    slide_architecture(
        prs, page=4, total=total, author_short=short,
        points=[
            "Гибридная схема ML + rules: ранкер ранжирует, composer гарантирует "
            "корректность (сплит на день/неделю, разнообразие, объём)",
            "Spec-driven разработка: spec 006 фиксирует требования к генератору, "
            "spec 012 — к ML-пайплайну",
            "Lazy-load ранкера с lru_cache; rule-based fallback при отсутствии "
            "или ошибке инференса — без видимого сбоя для пользователя",
            "Domain/-слой принимает ExercisePool dataclass'ы, возвращает PlannedPlan "
            "без I/O — изолированные unit-тесты на сложной логике composer'а",
            "Версионирование артефактов: ml/models/workout_rec/<algo>/v<semver>/",
        ],
        stack=[
            "Backend: FastAPI + SQLAlchemy 2 (async) + PostgreSQL 16",
            "ML: LightGBM, scikit-learn, pandas, joblib",
            "Frontend: Flutter (Web/PWA) + Riverpod + go_router + Dio",
            "Pipeline: ml/etl/workout_recommender/ + ml/training/workout_recommender/",
            "Инфраструктура: Docker Compose, Alembic-миграции, GitHub Actions",
            "Тесты: pytest + testcontainers PostgreSQL",
        ],
    )

    # 5. Алгоритм
    slide_content(
        prs, page=5, total=total, author_short=short,
        title="Гибридный алгоритм: ML-ранкер + composer",
        intro="Двухстадийная схема: ML сначала ранжирует, потом детерминированные "
              "правила собирают валидный план.",
        items=[
            "Стадия 1 (ML): LightGBM Classifier выдаёт score = P(упражнение "
            "подходит профилю); ранжирование каталога по этим score",
            "Прод-артефакт — один LightGBM-бустер на 320 деревьев и 36 признаков "
            "профиля + упражнения; бинарная классификация с целевой функцией "
            "log-loss и is_unbalance=True для несбалансированных меток",
            "Стадия 2 (composer): обход ранжированного каталога с проверкой "
            "сплита (push/pull/legs), оборудования, баланса compound/isolation",
            "Fallback path: при отсутствии ML-артефакта composer использует "
            "rule-based score (тот же интерфейс ExercisePool)",
            "Выход — иерархия workout_plan → plan_weeks → plan_days → plan_exercises, "
            "сохраняется в БД в одной транзакции",
        ],
    )

    # 6. Датасет
    slide_content(
        prs, page=6, total=total, author_short=short,
        title="Датасет: 97 760 пар user × exercise",
        intro="Открытых датасетов «профиль пользователя → подходящие упражнения» "
              "нет. Решение — программная разметка через rule-based labels.",
        items=[
            "Источник профилей: Kaggle gym-members-exercise-dataset "
            "(973 строки после фильтрации)",
            "Каталог упражнений: 100+ упражнений из API Ninjas + дополнения, "
            "разнесены по группам мышц / типу / оборудованию",
            "Разметка: для каждой пары (профиль, упражнение) применяются "
            "фиксированные правила → бинарная метка «рекомендуется/нет»",
            "Итог: 97 760 пар; сплит 70/15/15 без leakage по anon_user_id — "
            "68 432 train / 14 664 val / 14 664 test",
            "Воспроизводимость: SHA-256 датасета фиксируется в manifest.json; "
            "правила разметки версионированы вместе со скриптом ETL",
        ],
    )

    # 7. Сравнение моделей (chart)
    slide_chart(
        prs, page=7, total=total, author_short=short,
        title="Сравнение моделей ранжирования",
        chart=_chart_workout_models(),
        caption="LightGBM выбран в качестве основной модели: ROC-AUC=0.985, "
                "PR-AUC=0.993, P@8=0.99 на тестовом сплите 14 664 пар.",
    )

    # 8. Сравнительная таблица метрик
    slide_table(
        prs, page=8, total=total, author_short=short,
        title="Метрики моделей: ранжирование упражнений",
        headers=["Модель", "ROC-AUC", "PR-AUC", "Precision@8", "Recall@8"],
        rows=[
            ["Popularity baseline", "0,500", "0,679", "0,689", "0,104"],
            ["Logistic Regression", "0,939", "0,964", "0,952", "0,232"],
            ["LightGBM (GBDT)", "0,985", "0,993", "0,990", "0,241"],
            ["MLP (нейросеть)", "0,985", "0,993", "0,990", "0,241"],
        ],
        highlight_row=2,
        caption="LightGBM и MLP дают идентичную точность; LightGBM выбран основной "
                "за ×3–5 более быстрый инференс и нативный feature importance.",
    )

    # 9. Feature importance
    slide_chart(
        prs, page=9, total=total, author_short=short,
        title="Что внутри модели: feature importance",
        chart=_chart_feature_importance_egor(),
        caption="Главный сигнал — оборудование (user_equipment_count); это "
                "методологически осмысленно: рекомендовать упражнение со штангой "
                "пользователю без штанги не имеет смысла. Цель — на третьем месте.",
    )

    # 10. Composer
    slide_content(
        prs, page=10, total=total, author_short=short,
        title="Composer: детерминированная сборка плана",
        intro="Чистая функция в domain/-слое: принимает ExercisePool + ML-scores, "
              "возвращает PlannedPlan без обращений к БД.",
        items=[
            "Вход: каталог упражнений с rule-based score и опционально ML-scores; "
            "профиль пользователя; параметры цикла (4 недели по умолчанию)",
            "Шаг 1: сплит на дни недели по принципу push/pull/legs (или fullbody "
            "при частоте < 4 в неделю)",
            "Шаг 2: обход ранжированного каталога с гарантиями — баланс compound/isolation, "
            "разнообразие групп мышц, соответствие оборудованию",
            "Шаг 3: прогрессия — линейная по неделям цикла с deload-неделей "
            "при необходимости (REQ-08 spec 006)",
            "Изоляция от БД: composer тестируется чистыми unit-тестами без "
            "PostgreSQL, что даёт быстрый CI и простоту дебага",
        ],
    )

    # 11. Тестирование
    slide_content(
        prs, page=11, total=total, author_short=short,
        title="Тестирование",
        intro="Двухуровневая стратегия: domain/-юниты + интеграционные тесты "
              "плюс smoke-тест ML-пайплайна.",
        items=[
            "Composer: чистые unit-тесты без БД — проверки сплита, "
            "разнообразия, прогрессии, fallback'а",
            "Интеграционные тесты: testcontainers Postgres + savepoint per test, "
            "проверка end-to-end POST /api/v1/plans/generate",
            "ML-pipeline: smoke-тесты обучения popularity / LR / LightGBM на "
            "mini-датасете; проверка артефактов и manifest.json",
            "Тесты обработки граничных случаев: профиль не заполнен, оборудование "
            "не выбрано, ML-артефакт недоступен (fallback path)",
            "make check = lint (ruff) + typecheck (strict mypy на src/) + test; "
            "CI-проверка через GitHub Actions на каждый push",
        ],
    )

    # 12. UI (скриншоты)
    slide_screenshots(
        prs, page=12, total=total, author_short=short,
        title="Интерфейс приложения",
        shots=[
            (SCREENSHOTS / "screen-07.png",
             "Каталог упражнений с рекомендациями"),
            (SCREENSHOTS / "screen-08.png",
             "Активная тренировка по сгенерированному плану"),
            (SCREENSHOTS / "screen-04.png",
             "Обзор: статистика и история тренировок"),
        ],
    )

    # 13. Выводы
    slide_conclusion(
        prs, page=13, total=total, author_short=short,
        items=[
            "Реализовано кроссплатформенное PWA-приложение с end-to-end "
            "интеграцией ML-генератора планов в промышленную эксплуатацию",
            "Сформирован датасет 97 760 пар user × exercise с rule-based разметкой; "
            "обучены и сравнены 3 модели ранжирования",
            "LightGBM Classifier обеспечивает ROC-AUC=0.985, PR-AUC=0.993 — "
            "превосходит LR (0.939) и popularity-baseline (0.500)",
            "Реализован composer — детерминированный сборщик плана, "
            "гарантирующий корректность независимо от ML-скоров",
            "Архитектура спроектирована с расчётом на отказоустойчивость: "
            "rule-based fallback при отсутствии ML-артефакта",
        ],
    )

    # 12. Спасибо
    slide_thanks(prs, author="Ощепков Егор Сергеевич · группа МММ-401-О-03")

    out = HERE / "presentation-egor.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    maria = build_maria()
    egor = build_egor()
    print(f"Wrote: {maria}")
    print(f"Wrote: {egor}")
